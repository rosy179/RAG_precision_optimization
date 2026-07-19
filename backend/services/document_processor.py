"""
Multi-format document ingestion:
  PDF          → pdfplumber (+ GPT-4o Vision OCR fallback for scanned pages)
  DOCX         → python-docx (paragraphs + tables)
  PPTX         → python-pptx (slide text + tables + notes)
  XLSX         → openpyxl (sheets flattened to rows)
  TXT/Markdown → decode as text
  URL          → requests + BeautifulSoup
  Image        → OpenAI GPT-4o Vision
  Audio        → OpenAI Whisper transcription
"""

import os
import re
import sys
import json
import uuid
import base64
from pathlib import Path
from typing import Optional

_SRC = str(Path(__file__).parent.parent.parent / "src")
if _SRC not in sys.path:
    sys.path.insert(0, _SRC)

from dotenv import load_dotenv
load_dotenv()

from backend.services import llm

CHUNK_SIZE    = int(os.getenv("CHUNK_SIZE", 512))
CHUNK_OVERLAP = int(os.getenv("CHUNK_OVERLAP", 50))

# Contextual Retrieval (Anthropic-style): at ingest, an LLM writes a short
# context sentence per chunk ("this chunk covers X of document Y") that is
# prepended before embedding + BM25 — the ORIGINAL text is still what gets
# stored/displayed. Toggle off with CONTEXTUAL_RETRIEVAL=0.
CONTEXTUAL_RETRIEVAL   = os.getenv("CONTEXTUAL_RETRIEVAL", "1") not in ("0", "false", "False", "")
CONTEXT_MAX_DOC_WORDS  = 8000   # cap words of the doc sent for context generation
CONTEXT_EXCERPT_WORDS  = 40     # words of each chunk shown in the batch prompt

LLM_MODEL = os.getenv("LLM_MODEL", "gpt-4o-mini")

AUDIO_EXTENSIONS   = {".mp3", ".wav", ".m4a", ".ogg", ".webm"}
IMAGE_EXTENSIONS   = {".png", ".jpg", ".jpeg", ".webp"}
# Plain-text family (decoded directly) and Office XML family (parsed).
TEXT_EXTENSIONS    = {".txt", ".md", ".markdown"}
OFFICE_EXTENSIONS  = {".docx", ".pptx", ".xlsx"}
ALLOWED_EXTENSIONS = ({".pdf"} | TEXT_EXTENSIONS | OFFICE_EXTENSIONS
                      | IMAGE_EXTENSIONS | AUDIO_EXTENSIONS)


def _has_api_key() -> bool:
    return llm.has_api_key()


def _make_doc_id(name: str) -> str:
    # Must be unique per upload (not just per name) — two users, or the same
    # user re-uploading, can share a filename/URL, and Document.id is a
    # single global primary key.
    return uuid.uuid4().hex[:12]


_MD_HEADING_RE = re.compile(r'^(#{1,6})\s+(.+?)\s*#*$')


def _detect_heading(line: str) -> Optional[tuple[int, str]]:
    """Classify a line as a heading → (level, title), else None.

    Recognizes Markdown headings (also the '### Slide n' / '### SheetName'
    markers the DOCX/PPTX/XLSX handlers insert) and short ALL-CAPS lines.
    Kept deliberately conservative to avoid splitting prose into noise —
    heading-less documents simply fall back to plain word-window chunking.
    """
    s = line.strip()
    if len(s) < 3 or len(s) > 120:
        return None
    m = _MD_HEADING_RE.match(s)
    if m:
        return len(m.group(1)), m.group(2).strip()
    # ALL-CAPS short line (e.g. section banners in plain-text exports)
    letters = [c for c in s if c.isalpha()]
    if letters and len(s.split()) <= 10 and sum(c.isupper() for c in letters) / len(letters) > 0.85:
        return 2, s
    return None


def _page_at(word_idx: int, page_word_starts: Optional[list[int]]) -> Optional[int]:
    """1-based page whose first word index is the largest ≤ word_idx."""
    if not page_word_starts:
        return None
    page = 1
    for p, start in enumerate(page_word_starts, 1):
        if word_idx >= start:
            page = p
        else:
            break
    return page


def _chunk_document(text: str, doc_id: str, title: str, source_type: str,
                    page_word_starts: Optional[list[int]] = None) -> list[dict]:
    """Heading-aware overlapping word-window chunker.

    Walks the text line by line, tracking the current heading path (a stack
    of headings by level). Chunks never straddle a new heading — a section
    boundary flushes the current window and the next chunk starts clean (no
    cross-section overlap) — and each chunk is tagged with the `heading` path
    active at its first word. Within a long section the usual sliding window
    with CHUNK_OVERLAP applies. Falls back to a single flat window when the
    document has no detectable headings.

    page_word_starts: word index where each page begins (PDFs) — used to tag
    every chunk with the 1-based page its first word falls on.
    """
    words: list[str] = []
    word_heading: list[str] = []      # heading path parallel to `words`
    boundaries: set[int] = set()      # word indices where a new section starts
    path: list[tuple[int, str]] = []  # stack of (level, title)

    for line in text.split("\n"):
        h = _detect_heading(line)
        if h:
            level, htitle = h
            while path and path[-1][0] >= level:
                path.pop()
            path.append((level, htitle))
            boundaries.add(len(words))
        cur = " > ".join(t for _, t in path)
        for w in line.split():
            words.append(w)
            word_heading.append(cur)

    if not words:
        return []

    chunks: list[dict] = []
    i = idx = 0
    step = max(1, CHUNK_SIZE - CHUNK_OVERLAP)
    while i < len(words):
        end = min(i + CHUNK_SIZE, len(words))
        # Don't cross into a new section: cut at the next heading boundary.
        next_bound = min((b for b in boundaries if i < b < end), default=None)
        cut_at_boundary = next_bound is not None
        if cut_at_boundary:
            end = next_bound
        chunk = {
            "id":     f"{doc_id}_c{idx}",
            "text":   " ".join(words[i:end]),
            "title":  title,
            "source": source_type,
            "doc_id": doc_id,
            "heading": word_heading[i],
        }
        page = _page_at(i, page_word_starts)
        if page is not None:
            chunk["page"] = page
        chunks.append(chunk)
        idx += 1
        # Start the next chunk exactly at a section boundary (no overlap
        # bleeding across sections); otherwise slide with overlap.
        i = end if cut_at_boundary else (end if end >= len(words) else i + step)
    return chunks


# ── Contextual Retrieval (Anthropic-style) ────────────────

def contextual_text(chunk: dict) -> str:
    """Text used for EMBEDDING and BM25: the generated context sentence
    prepended to the original chunk text. The original `text` is what stays
    stored/displayed — this is only the retrieval representation."""
    ctx = chunk.get("context")
    return f"{ctx}\n\n{chunk['text']}" if ctx else chunk["text"]


def _fallback_context(title: str, heading: str) -> str:
    """Deterministic context when the LLM is unavailable or misaligns —
    still situates the chunk by document + section for retrieval."""
    if heading:
        return f"Trích từ mục \"{heading}\" của tài liệu \"{title}\"."
    return f"Trích từ tài liệu \"{title}\"."


# English instruction (kept language-neutral on purpose): a Vietnamese system
# prompt biased gpt-4o-mini into writing Vietnamese context for English
# documents, mixing languages inside one embedding. The context must be in the
# document's OWN language so it aligns cleanly with same-language chunk text.
_CONTEXT_SYSTEM_PROMPT = (
    "You situate each text chunk within its source document to improve search "
    "retrieval. Reply strictly in the requested JSON format."
)


def _generate_contexts(title: str, full_text: str, chunks: list[dict]) -> Optional[list[str]]:
    """One LLM call situating every chunk of a document (Anthropic Contextual
    Retrieval, batched per document to keep cost/latency low — the document is
    sent once, each chunk only as a short excerpt). Returns a context per
    chunk, or None on any failure/misalignment so the caller can fall back."""
    if not _has_api_key():
        return None
    doc = " ".join(full_text.split()[:CONTEXT_MAX_DOC_WORDS])
    listing = "\n".join(
        f"[{i}] " + (f"(section: {c['heading']}) " if c.get("heading") else "")
        + " ".join(c["text"].split()[:CONTEXT_EXCERPT_WORDS])
        for i, c in enumerate(chunks, 1)
    )
    n = len(chunks)
    user = (
        f"<document title=\"{title}\">\n{doc}\n</document>\n\n"
        f"The document above is split into {n} chunks. For EACH chunk, write ONE "
        f"short sentence (max 25 words) stating which section/topic of the "
        f"document it covers and its key entities/concepts, to improve search. "
        f"Write each context in the SAME LANGUAGE as the document text.\n\n"
        f"Chunks:\n{listing}\n\n"
        f"Return JSON: {{\"contexts\": [\"...\"]}} with EXACTLY {n} items in order."
    )
    try:
        client = llm.get_client()
        resp = client.chat.completions.create(
            model=LLM_MODEL,
            messages=[{"role": "system", "content": _CONTEXT_SYSTEM_PROMPT},
                      {"role": "user", "content": user}],
            temperature=0.0,
            max_tokens=min(2000, 80 + 45 * n),
            response_format={"type": "json_object"},
        )
        arr = json.loads(resp.choices[0].message.content or "{}").get("contexts")
        if isinstance(arr, list) and len(arr) == n:
            return [str(x).strip() for x in arr]
    except Exception:
        pass
    return None


def _apply_contextual_retrieval(title: str, full_text: str, chunks: list[dict]) -> None:
    """Attach a `context` sentence to each chunk (in place). Skipped when the
    feature is off or the document is a single chunk (nothing to situate)."""
    if not CONTEXTUAL_RETRIEVAL or len(chunks) <= 1:
        return
    generated = _generate_contexts(title, full_text, chunks)
    for i, c in enumerate(chunks):
        ctx = generated[i] if generated and generated[i] else _fallback_context(title, c.get("heading", ""))
        if ctx:
            c["context"] = ctx


def _finalize_text_doc(full_text: str, doc_id: str, title: str, source_type: str,
                       page_word_starts: Optional[list[int]] = None) -> tuple[list[dict], str]:
    """Shared tail for text-bearing documents: heading-aware chunking →
    contextual retrieval → hierarchical summary. Returns (chunks, summary)."""
    chunks = _chunk_document(full_text, doc_id, title, source_type, page_word_starts)
    _apply_contextual_retrieval(title, full_text, chunks)
    summary = _generate_summary(full_text[:4000], title)
    return chunks, summary


def reconstruct_from_chunks(chunk_texts: list[str]) -> tuple[str, list[tuple[int, int]]]:
    """Rebuild a document's canonical text from its ordered sliding-window
    chunks, returning (text, [(char_start, char_end) per chunk]).

    The overlap is found by suffix/prefix word matching rather than trusting
    CHUNK_OVERLAP, so documents ingested under older chunking settings still
    reconstruct correctly.
    """
    words: list[str] = []
    ranges: list[tuple[int, int]] = []  # word ranges per chunk
    for text in chunk_texts:
        cw = text.split()
        overlap = 0
        max_k = min(len(words), len(cw))
        for k in range(max_k, 0, -1):
            if words[len(words) - k:] == cw[:k]:
                overlap = k
                break
        start_word = len(words) - overlap
        words.extend(cw[overlap:])
        ranges.append((start_word, start_word + len(cw)))

    # word index → char offset in " ".join(words)
    char_at: list[int] = []
    pos = 0
    for w in words:
        char_at.append(pos)
        pos += len(w) + 1
    text = " ".join(words)
    offsets = [
        (char_at[s] if s < len(char_at) else len(text),
         (char_at[e - 1] + len(words[e - 1])) if 0 < e <= len(words) else len(text))
        for s, e in ranges
    ]
    return text, offsets


def _clean_text(text: str) -> str:
    """Remove excessive whitespace/newlines."""
    text = re.sub(r'\n{3,}', '\n\n', text)
    text = re.sub(r' {2,}', ' ', text)
    return text.strip()


# ── Format Handlers ───────────────────────────────────────

# A page whose extractable text is shorter than this is treated as scanned
# (image-only) and sent through OCR — pdfplumber returns "" or a stray
# header for such pages, which used to index as 0 useful chunks silently.
_MIN_PAGE_TEXT_CHARS = 40


def process_pdf(file_bytes: bytes, filename: str) -> tuple[list[dict], str]:
    """
    Returns (chunks, doc_summary).

    Pages with a real text layer are read by pdfplumber; pages with (almost)
    no extractable text are assumed scanned and OCR'd via GPT-4o Vision, so a
    scanned PDF no longer indexes as zero chunks in silence.
    """
    import pdfplumber, io
    # Per-page text ("" for pages with no usable text layer → OCR candidates).
    page_texts: list[str] = []
    with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
        for page in pdf.pages:
            t = page.extract_text() or ""
            page_texts.append(_clean_text(t) if len(t.strip()) >= _MIN_PAGE_TEXT_CHARS else "")

    # Splice OCR text into the image-only pages, preserving page order.
    scanned_pages = [i for i, t in enumerate(page_texts) if not t]
    if scanned_pages and _has_api_key():
        for pnum, text in _ocr_pdf_pages(file_bytes, scanned_pages, filename).items():
            page_texts[pnum] = text

    # Tag each chunk with the 1-based page its first word starts on, so the
    # viewer can jump straight there. Headings live line by line, so keep the
    # newlines between pages for the heading-aware chunker to see.
    page_word_starts: list[int] = []
    words: list[str] = []
    for t in page_texts:
        page_word_starts.append(len(words))
        words.extend(t.split())

    full_text = "\n".join(t for t in page_texts if t)
    doc_id = _make_doc_id(filename)
    return _finalize_text_doc(full_text, doc_id, filename, "pdf",
                              page_word_starts=page_word_starts)


def _ocr_pdf_pages(file_bytes: bytes, page_indices: list[int],
                   filename: str) -> dict[int, str]:
    """Render the given 0-based PDF pages to images (PyMuPDF, 2x zoom for
    legibility) and OCR each via GPT-4o Vision. Returns {page_index: text}
    for pages that yielded text. Best-effort: a failed page is skipped."""
    import fitz  # PyMuPDF

    out: dict[int, str] = {}
    try:
        doc = fitz.open(stream=file_bytes, filetype="pdf")
    except Exception:
        return out
    # Cap the OCR work so a huge scanned PDF can't run up an unbounded bill.
    wanted = set(page_indices[:30])
    try:
        zoom = fitz.Matrix(2, 2)
        for pnum in range(doc.page_count):
            if pnum not in wanted:
                continue
            try:
                png = doc.load_page(pnum).get_pixmap(matrix=zoom).tobytes("png")
                # The OCR prompt returns "" for a blank page, so any non-empty
                # transcription is real text (unlike the text-layer heuristic,
                # this is not gated on _MIN_PAGE_TEXT_CHARS).
                text = _clean_text(_vision_extract_text(png, "image/png"))
                if text:
                    out[pnum] = text
            except Exception:
                continue
    finally:
        doc.close()
    return out


# URL ingestion download cap — bounds both memory use and what an SSRF
# attempt could exfiltrate in one response.
MAX_FETCH_BYTES = 5 * 1024 * 1024


def _require_public_url(url: str) -> None:
    """Reject URLs that would make the server request itself or the internal
    network (SSRF): non-http(s) schemes, and hosts resolving to loopback,
    private, link-local or any other non-global address (this also covers
    cloud metadata endpoints like 169.254.169.254)."""
    import ipaddress
    import socket
    from urllib.parse import urlparse

    parsed = urlparse(url)
    if parsed.scheme not in ("http", "https"):
        raise ValueError("Chỉ hỗ trợ URL http/https")
    if not parsed.hostname:
        raise ValueError("URL không hợp lệ")
    try:
        infos = socket.getaddrinfo(
            parsed.hostname, parsed.port or (443 if parsed.scheme == "https" else 80),
            proto=socket.IPPROTO_TCP)
    except socket.gaierror as e:
        raise ValueError(f"Không phân giải được tên miền: {parsed.hostname}") from e
    for info in infos:
        ip = ipaddress.ip_address(info[4][0])
        if not ip.is_global:
            raise ValueError(
                "URL trỏ tới địa chỉ mạng nội bộ — đã chặn để tránh SSRF")


def _fetch_public_url(url: str) -> tuple[bytes, str]:
    """GET with SSRF validation re-run on every redirect hop (max 5) and a
    MAX_FETCH_BYTES download cap. Returns (body, encoding guess).

    Known limit: validation resolves DNS separately from the request itself,
    so a hostile nameserver flipping records between the two lookups (DNS
    rebinding) is not covered — acceptable for this deployment.
    """
    import requests
    from urllib.parse import urljoin

    headers = {
        "User-Agent": ("Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 "
                       "(KHTML, like Gecko) Chrome/126.0 Safari/537.36"),
        "Accept-Language": "en-US,en;q=0.9,vi;q=0.8",
    }
    for _ in range(5):
        _require_public_url(url)
        resp = requests.get(url, timeout=20, headers=headers,
                            allow_redirects=False, stream=True)
        if resp.status_code in (301, 302, 303, 307, 308):
            loc = resp.headers.get("Location")
            if not loc:
                raise ValueError("Chuyển hướng không có địa chỉ đích")
            url = urljoin(url, loc)
            continue
        resp.raise_for_status()
        raw = bytearray()
        for chunk in resp.iter_content(65536):
            raw.extend(chunk)
            if len(raw) > MAX_FETCH_BYTES:
                raise ValueError("Trang quá lớn (giới hạn 5 MB)")
        # requests defaults to ISO-8859-1 when the header has no charset —
        # sniff the meta tag instead, falling back to UTF-8.
        enc = resp.encoding
        if not enc or enc.lower() == "iso-8859-1":
            m = re.search(rb'charset=["\']?([\w-]{2,20})', bytes(raw[:4096]), re.I)
            enc = m.group(1).decode("ascii", "ignore") if m else "utf-8"
        return bytes(raw), enc
    raise ValueError("Quá nhiều lần chuyển hướng (redirect)")


def process_url(url: str) -> tuple[list[dict], str]:
    raw, enc = _fetch_public_url(url)
    try:
        html = raw.decode(enc, errors="replace")
    except LookupError:
        html = raw.decode("utf-8", errors="replace")

    text, title = "", ""

    # Primary: trafilatura locates the actual article body. Naive
    # first-<main>/<article> selection grabs cookie banners and promo
    # widgets on real-world pages (that's how we indexed 2 KB of cookie
    # consent text instead of a 36 KB article).
    try:
        import trafilatura
        text = trafilatura.extract(html, include_comments=False, include_tables=True) or ""
        meta = trafilatura.extract_metadata(html)
        if meta and meta.title:
            title = meta.title
    except Exception:
        pass

    # Fallback: BeautifulSoup with junk stripping, picking the candidate
    # with the MOST text instead of the first one in document order.
    if len(text) < 200:
        from bs4 import BeautifulSoup
        soup = BeautifulSoup(html, "html.parser")
        if not title and soup.title and soup.title.string:
            title = soup.title.string.strip()
        for tag in soup(["script", "style", "nav", "footer", "header", "aside",
                         "noscript", "form", "iframe", "svg", "button"]):
            tag.decompose()
        junk = re.compile(r"cookie|consent|banner|popup|modal|menu|sidebar|share|"
                          r"social|subscribe|newsletter|advert|promo|breadcrumb", re.I)
        for attr in ("class", "id"):
            for tag in soup.find_all(attrs={attr: junk}):
                tag.decompose()
        candidates = soup.find_all(["main", "article"]) or ([soup.body] if soup.body else [])
        best = max(candidates, key=lambda t: len(t.get_text()), default=None)
        text = (best or soup).get_text(separator="\n", strip=True)

    if not title:
        m = re.search(r"<title[^>]*>(.*?)</title>", html, re.I | re.S)
        title = _clean_text(m.group(1)) if m else url

    text = _clean_text(text)
    if len(text) < 100:
        raise ValueError("Trang không có nội dung văn bản đọc được (có thể nội dung được render bằng JavaScript)")

    doc_id = _make_doc_id(url)
    return _finalize_text_doc(text, doc_id, title, "url")


_VISION_DESCRIBE_PROMPT = (
    "You are a technical knowledge base assistant. Describe this image/diagram "
    "in detail for indexing purposes. Include: all text labels, values, trends, "
    "components, relationships, and key technical insights. Be thorough."
)
_VISION_OCR_PROMPT = (
    "Transcribe ALL text visible in this document page exactly, preserving "
    "reading order, headings, lists and table structure. Output only the "
    "transcribed text, no commentary. If the page has no readable text, output "
    "nothing."
)


def _vision_call(image_bytes: bytes, mime: str, prompt: str,
                 max_tokens: int = 1000) -> str:
    """One GPT-4o Vision call over an inline base64 image. Callers must
    ensure an API key exists (_has_api_key)."""
    b64 = base64.b64encode(image_bytes).decode()
    client = llm.get_client()
    resp = client.chat.completions.create(
        model="gpt-4o",
        messages=[{"role": "user", "content": [
            {"type": "image_url", "image_url": {"url": f"data:{mime};base64,{b64}"}},
            {"type": "text", "text": prompt},
        ]}],
        max_tokens=max_tokens,
    )
    return (resp.choices[0].message.content or "").strip()


def _vision_extract_text(image_bytes: bytes, mime: str) -> str:
    """OCR a rendered page image into plain text (scanned-PDF fallback)."""
    return _vision_call(image_bytes, mime, _VISION_OCR_PROMPT, max_tokens=1500)


def process_image(file_bytes: bytes, filename: str) -> tuple[list[dict], str]:
    if not _has_api_key():
        description = f"[MOCK] Image description for {filename}"
    else:
        ext = Path(filename).suffix.lower().lstrip(".")
        mime = "image/png" if ext == "png" else "image/jpeg"
        description = _vision_call(file_bytes, mime, _VISION_DESCRIBE_PROMPT)

    doc_id = _make_doc_id(filename)
    # A vision description is already a whole-document blob — no heading
    # structure and nothing to situate, so no contextual retrieval here.
    chunks = _chunk_document(description, doc_id, filename, "image")
    summary = description[:500]
    return chunks, summary


def transcribe_audio(file_bytes: bytes, filename: str) -> str:
    """Whisper speech-to-text; mock text when no API key is configured."""
    if llm.is_mock():
        return f"[MOCK] Audio transcription for {filename}"
    client = llm.get_client()
    resp = client.audio.transcriptions.create(
        model="whisper-1",
        file=(filename, file_bytes),
    )
    return resp.text.strip()


def process_audio(file_bytes: bytes, filename: str) -> tuple[list[dict], str]:
    transcript = _clean_text(transcribe_audio(file_bytes, filename))
    doc_id = _make_doc_id(filename)
    return _finalize_text_doc(transcript, doc_id, filename, "audio")


def process_txt(file_bytes: bytes, filename: str,
                source_type: str = "txt") -> tuple[list[dict], str]:
    text = _clean_text(file_bytes.decode("utf-8", errors="replace"))
    doc_id = _make_doc_id(filename)
    return _finalize_text_doc(text, doc_id, filename, source_type)


def process_docx(file_bytes: bytes, filename: str) -> tuple[list[dict], str]:
    """Word .docx → paragraphs + table cells, in document order."""
    import docx, io

    document = docx.Document(io.BytesIO(file_bytes))
    parts: list[str] = []
    for para in document.paragraphs:
        if para.text.strip():
            parts.append(para.text)
    for table in document.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                parts.append(" | ".join(cells))

    text = _clean_text("\n".join(parts))
    doc_id = _make_doc_id(filename)
    return _finalize_text_doc(text, doc_id, filename, "docx")


def process_pptx(file_bytes: bytes, filename: str) -> tuple[list[dict], str]:
    """PowerPoint .pptx → per-slice text frames, tables and speaker notes."""
    import pptx, io

    presentation = pptx.Presentation(io.BytesIO(file_bytes))
    parts: list[str] = []
    for i, slide in enumerate(presentation.slides, 1):
        parts.append(f"### Slide {i}")
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                parts.append(shape.text_frame.text)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells if c.text.strip()]
                    if cells:
                        parts.append(" | ".join(cells))
        notes = slide.notes_slide.notes_text_frame.text if slide.has_notes_slide else ""
        if notes and notes.strip():
            parts.append(f"(Ghi chú: {notes.strip()})")

    text = _clean_text("\n".join(parts))
    doc_id = _make_doc_id(filename)
    return _finalize_text_doc(text, doc_id, filename, "pptx")


def process_xlsx(file_bytes: bytes, filename: str) -> tuple[list[dict], str]:
    """Excel .xlsx → each sheet flattened to 'col | col | col' rows."""
    import openpyxl, io

    wb = openpyxl.load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
    parts: list[str] = []
    try:
        for ws in wb.worksheets:
            parts.append(f"### {ws.title}")
            for row in ws.iter_rows(values_only=True):
                cells = ["" if v is None else str(v) for v in row]
                if any(c.strip() for c in cells):
                    parts.append(" | ".join(cells).rstrip(" |"))
    finally:
        wb.close()

    text = _clean_text("\n".join(parts))
    doc_id = _make_doc_id(filename)
    return _finalize_text_doc(text, doc_id, filename, "xlsx")


def process_file(file_bytes: bytes, filename: str,
                 content_type: str = "") -> tuple[list[dict], str, str]:
    """Dispatch to the right handler by extension/MIME.

    Returns (chunks, summary, doc_type). Raises ValueError for
    unsupported formats.
    """
    ext = "." + filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
    if ext == ".pdf" or "pdf" in content_type:
        chunks, summary = process_pdf(file_bytes, filename)
        return chunks, summary, "pdf"
    if ext in AUDIO_EXTENSIONS or content_type.startswith("audio/"):
        chunks, summary = process_audio(file_bytes, filename)
        return chunks, summary, "audio"
    if ext in IMAGE_EXTENSIONS or content_type.startswith("image/"):
        chunks, summary = process_image(file_bytes, filename)
        return chunks, summary, "image"
    if ext == ".docx":
        chunks, summary = process_docx(file_bytes, filename)
        return chunks, summary, "docx"
    if ext == ".pptx":
        chunks, summary = process_pptx(file_bytes, filename)
        return chunks, summary, "pptx"
    if ext == ".xlsx":
        chunks, summary = process_xlsx(file_bytes, filename)
        return chunks, summary, "xlsx"
    # Markdown/txt and any remaining text/* payloads decode as plain text.
    if ext in TEXT_EXTENSIONS or content_type.startswith("text/"):
        source = "markdown" if ext in (".md", ".markdown") else "txt"
        chunks, summary = process_txt(file_bytes, filename, source)
        return chunks, summary, source
    raise ValueError(f"Unsupported file type: {ext or content_type or 'unknown'}")


# ── Document Summary ──────────────────────────────────────

def _generate_summary(text: str, title: str) -> str:
    """Generate a short summary for hierarchical indexing."""
    if llm.is_mock():
        return text[:300]
    client = llm.get_client()
    try:
        resp = client.chat.completions.create(
            model=LLM_MODEL,
            messages=[{"role": "user", "content": (
                f"Summarize this document in 3-5 sentences for a technical search index. "
                f"Document title: {title}\n\n{text}"
            )}],
            max_tokens=200,
            temperature=0.0,
        )
        return resp.choices[0].message.content.strip()
    except Exception:
        return text[:300]
